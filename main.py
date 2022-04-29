import os
import shutil
import datetime
import re
import zipfile
from pptx import Presentation
from tika import parser
from tkinter import *
from tkinter.ttk import *

my_lectures = ['3D 게임 프로그래밍', 'STL', '네트워크 기초', '선형대수학', '스크립트 언어', '인간과 철학']
normal_folders = ['docx', 'pptx', 'pdf', 'hwp', 'jpg', 'png', 'txt', 'zip']

l3d = re.compile(r'(3d|3D)?((game|Game)[ ]?(Programming|programming)|게임[ ]?프로그래밍)(DirectX|Direct3D)?')
lstl = re.compile(r'stl|STL')
lnw = re.compile(r'network|네트워크[ ]?(기초)?|Network')
lln = re.compile(r'linear|선형(대)?[ ]?수학?|Linear')
lsc = re.compile(r'(script|스크립트[ ]?언어|Script)(Python|python|pycharm|Pycharm)?(.py)?')
lhp = re.compile(r'인간과[ ]?철학|인철|철학')


def make_lecture_folders():
    if os.path.exists('/Lab/Univ/3-1'):
        os.chdir('/Lab/Univ/3-1')
    else:
        os.makedirs('/Lab/Univ/3-1')
        os.chdir('/Lab/Univ/3-1')

    for lecture in my_lectures:
        if os.path.exists(lecture):
            continue
        os.mkdir(lecture)
    print("done")


def make_normal_folders():
    if os.path.exists('/Lab/Files'):
        os.chdir('/Lab/Files')
    else:
        os.makedirs('/Lab/Files')
        os.chdir('/Lab/Files')

    for category in normal_folders:
        if os.path.exists(category):
            continue
        os.mkdir(category)
    print("done")


l3d = re.compile(r'(3d|3D)?((game|Game)[ ]?(Programming|programming)|게임[ ]?프로그래밍)?(DirectX|Direct3D)')
lstl = re.compile(r'(stl|STL)(vector)?')
lnw = re.compile(r'network|네트워크[ ]?(기초)?|Network')
lln = re.compile(r'linear|선형(대)?[ ]?(수학)?|Linear')
lsc = re.compile(r'(script|스크립트[ ]?(언어)?|Script|Python|python|pycharm|Pycharm)(.py)?')
lhp = re.compile(r'인간과[ ]?철학|인철|철학')

def move_files():
    os.chdir('/Lab/Downloads')
    now = datetime.datetime.today().strftime('%y%m%d')

    file_list = []
    # 오늘 다운받은 파일만 해당
    for file in os.listdir():
        abs_path = os.path.abspath(file)
        ctime = os.path.getctime(abs_path)
        cdate = datetime.datetime.fromtimestamp(ctime).strftime('%y%m%d')
        if cdate == now:
            file_list.append(file)

    for file in file_list:
        result = []
        result.append(file)

        # if file.endswith('.zip'):
        #     fn = file.split('.zip')
        #     fn = fn[0]
        #     zf = zipfile.ZipFile(file)
        #     zf.extractall('./' + fn)
        #     zf.close()
        #     os.chdir('./' + fn)
        if file.endswith('.pdf'):
            data = parser.from_file(file)
            content = data['content'].strip()
            result.append(content)
        elif file.endswith('.pptx'):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        result.append(paragraph.text)
        elif file.endswith('.txt'):
            try:
                f = open(file, encoding='UTF-8')
                content = f.read()
                result.append(content)
                f.close()
            except Exception as e:
                print(e)

        for word in result:
            if lstl.search(word):
                # shutil.move(file, '/Lab/Univ/3-1/STL')
                shutil.copy(file, '/Lab/Univ/3-1/STL')
            elif l3d.search(word):
                # shutil.move(file, '/Lab/Univ/3-1/3d 게임 프로그래밍')
                shutil.copy(file, '/Lab/Univ/3-1/3d 게임 프로그래밍')
            elif lnw.search(word):
                # shutil.move(file, '/Lab/Univ/3-1/네트워크 기초')
                shutil.copy(file, '/Lab/Univ/3-1/네트워크 기초')
            elif lln.search(word):
                # shutil.move(file, '/Lab/Univ/3-1/선형대수학')
                shutil.copy(file, '/Lab/Univ/3-1/선형대수학')
            elif lsc.search(word):
                # shutil.move(file, '/Lab/Univ/3-1/스크립트 언어')
                shutil.copy(file, '/Lab/Univ/3-1/스크립트 언어')
            elif lhp.search(word):
                # shutil.move(file, '/Lab/Univ/3-1/인간과 철학')
                shutil.copy(file, '/Lab/Univ/3-1/인간과 철학')
            else:
                for n in normal_folders:
                    if file.endswith('.' + n):
                        shutil.copy(file, '/Lab/Files/' + n)

    print("done")


def stop():
    window.quit()


# shutil.rmtree('/Lab/Downloads')
# shutil.copytree('/Lab/backup', '/Lab/Downloads')
if os.path.exists('/Lab/Files'):
    shutil.rmtree('/Lab/Files')
if os.path.exists('/Lab/Univ'):
    shutil.rmtree('/Lab/Univ')


window = Tk()
window.geometry("+0+0")
window.resizable(False, False)
window.bind('<Escape>', stop)

label = Label(text="File Management")
label.pack()
# button = Button(text="Start", takefocus=False)
# button.pack()
command_frame = LabelFrame(text='Command')
command_frame.pack(fill=BOTH, padx=5, pady=5)

Button(command_frame, command=make_lecture_folders, text='make lecture_folders').pack(side=LEFT, expand=True, fill=BOTH, padx=5, pady=5)
Button(command_frame, command=make_normal_folders, text='make normal_folders').pack(side=LEFT, expand=True, fill=BOTH, padx=5, pady=5)
Button(command_frame, command=move_files, text='move files').pack(side=LEFT, expand=True, fill=BOTH, padx=5, pady=5)

menu = Menu()
menu_File = Menu(menu, tearoff=False) # tearoff : menu 분리
menu_File.add_command(label="Quit", accelerator='Ctrl+Q', command=stop)
menu.add_cascade(label="Menu", underline=True, menu=menu_File)
menu_Colors = Menu(menu, tearoff=False)
window.bind_all('<Control-q>', stop)
window.config(menu=menu)

window.mainloop()